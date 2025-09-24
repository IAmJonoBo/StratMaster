"""Document processing pipeline for multiple formats (docx, pdf, pptx, md)."""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any, Dict, List, Optional, Union

import markdown
import nltk
import spacy
from docx import Document as DocxDocument
from pptx import Presentation
from PyPDF2 import PdfReader
from pydantic import BaseModel, Field

logger = logging.getLogger(__name__)

# Download required NLTK data
try:
    nltk.data.find('tokenizers/punkt')
except LookupError:
    nltk.download('punkt')

try:
    nltk.data.find('corpora/stopwords')
except LookupError:
    nltk.download('stopwords')


class ExtractedEntity(BaseModel):
    """Extracted entity from document content."""
    text: str
    label: str
    confidence: float
    start_char: int
    end_char: int


class DocumentSection(BaseModel):
    """A section of processed document content."""
    title: str
    content: str
    section_type: str = "content"  # content, header, bullet, table
    page_number: Optional[int] = None
    slide_number: Optional[int] = None
    entities: List[ExtractedEntity] = Field(default_factory=list)


class ProcessedDocument(BaseModel):
    """Complete processed document with extracted content and metadata."""
    filename: str
    file_type: str
    total_pages: Optional[int] = None
    total_slides: Optional[int] = None
    sections: List[DocumentSection] = Field(default_factory=list)
    entities: List[ExtractedEntity] = Field(default_factory=list)
    key_facts: List[str] = Field(default_factory=list)
    summary: str = ""
    metadata: Dict[str, Any] = Field(default_factory=dict)


class DocumentProcessor:
    """Processes multiple document formats and extracts structured content."""
    
    def __init__(self, spacy_model: str = "en_core_web_sm"):
        """Initialize document processor with NLP model.
        
        Args:
            spacy_model: SpaCy model name for entity extraction
        """
        try:
            self.nlp = spacy.load(spacy_model)
        except OSError:
            logger.warning(f"SpaCy model {spacy_model} not found. Install with: python -m spacy download {spacy_model}")
            self.nlp = None
        
        self.stop_words = set(nltk.corpus.stopwords.words('english'))
    
    def process_file(self, file_path: Union[str, Path]) -> ProcessedDocument:
        """Process a file and return structured content.
        
        Args:
            file_path: Path to the document file
            
        Returns:
            ProcessedDocument with extracted content and metadata
        """
        file_path = Path(file_path)
        file_type = file_path.suffix.lower()
        
        logger.info(f"Processing {file_type} file: {file_path.name}")
        
        if file_type == '.docx':
            return self._process_docx(file_path)
        elif file_type == '.pdf':
            return self._process_pdf(file_path)
        elif file_type in ['.pptx', '.ppt']:
            return self._process_pptx(file_path)
        elif file_type in ['.md', '.markdown']:
            return self._process_markdown(file_path)
        else:
            raise ValueError(f"Unsupported file type: {file_type}")
    
    def _process_docx(self, file_path: Path) -> ProcessedDocument:
        """Process Microsoft Word document."""
        doc = DocxDocument(file_path)
        sections = []
        
        for i, paragraph in enumerate(doc.paragraphs):
            if paragraph.text.strip():
                section_type = "header" if paragraph.style.name.startswith('Heading') else "content"
                
                section = DocumentSection(
                    title=f"Paragraph {i+1}" if section_type == "content" else paragraph.text[:50],
                    content=paragraph.text,
                    section_type=section_type
                )
                
                # Extract entities from content
                if self.nlp:
                    entities = self._extract_entities(paragraph.text)
                    section.entities = entities
                
                sections.append(section)
        
        # Process tables
        for i, table in enumerate(doc.tables):
            table_content = self._extract_table_text(table)
            if table_content:
                section = DocumentSection(
                    title=f"Table {i+1}",
                    content=table_content,
                    section_type="table"
                )
                sections.append(section)
        
        return self._create_processed_document(
            filename=file_path.name,
            file_type="docx",
            sections=sections
        )
    
    def _process_pdf(self, file_path: Path) -> ProcessedDocument:
        """Process PDF document."""
        sections = []
        
        with open(file_path, 'rb') as file:
            pdf_reader = PdfReader(file)
            total_pages = len(pdf_reader.pages)
            
            for page_num, page in enumerate(pdf_reader.pages, 1):
                text = page.extract_text()
                if text.strip():
                    # Split into paragraphs
                    paragraphs = [p.strip() for p in text.split('\n\n') if p.strip()]
                    
                    for i, paragraph in enumerate(paragraphs):
                        section = DocumentSection(
                            title=f"Page {page_num}, Paragraph {i+1}",
                            content=paragraph,
                            section_type="content",
                            page_number=page_num
                        )
                        
                        # Extract entities
                        if self.nlp:
                            entities = self._extract_entities(paragraph)
                            section.entities = entities
                        
                        sections.append(section)
        
        return self._create_processed_document(
            filename=file_path.name,
            file_type="pdf",
            sections=sections,
            total_pages=total_pages
        )
    
    def _process_pptx(self, file_path: Path) -> ProcessedDocument:
        """Process PowerPoint presentation."""
        prs = Presentation(file_path)
        sections = []
        
        for slide_num, slide in enumerate(prs.slides, 1):
            slide_content = []
            
            # Extract text from shapes
            for shape in slide.shapes:
                if hasattr(shape, "text") and shape.text.strip():
                    slide_content.append(shape.text.strip())
            
            if slide_content:
                content = "\n".join(slide_content)
                section = DocumentSection(
                    title=f"Slide {slide_num}",
                    content=content,
                    section_type="slide",
                    slide_number=slide_num
                )
                
                # Extract entities
                if self.nlp:
                    entities = self._extract_entities(content)
                    section.entities = entities
                
                sections.append(section)
        
        return self._create_processed_document(
            filename=file_path.name,
            file_type="pptx",
            sections=sections,
            total_slides=len(prs.slides)
        )
    
    def _process_markdown(self, file_path: Path) -> ProcessedDocument:
        """Process Markdown document."""
        with open(file_path, 'r', encoding='utf-8') as file:
            content = file.read()
        
        # Parse markdown to HTML then extract structured content
        html = markdown.markdown(content, extensions=['toc', 'tables', 'fenced_code'])
        
        # Split by headers for sections
        sections = []
        header_pattern = r'^(#{1,6})\s+(.+)$'
        current_section = []
        current_title = "Introduction"
        
        for line in content.split('\n'):
            header_match = re.match(header_pattern, line)
            if header_match:
                # Save previous section
                if current_section:
                    section_content = '\n'.join(current_section).strip()
                    if section_content:
                        section = DocumentSection(
                            title=current_title,
                            content=section_content,
                            section_type="header"
                        )
                        
                        # Extract entities
                        if self.nlp:
                            entities = self._extract_entities(section_content)
                            section.entities = entities
                        
                        sections.append(section)
                
                # Start new section
                current_title = header_match.group(2)
                current_section = []
            else:
                current_section.append(line)
        
        # Don't forget the last section
        if current_section:
            section_content = '\n'.join(current_section).strip()
            if section_content:
                section = DocumentSection(
                    title=current_title,
                    content=section_content,
                    section_type="content"
                )
                
                if self.nlp:
                    entities = self._extract_entities(section_content)
                    section.entities = entities
                
                sections.append(section)
        
        return self._create_processed_document(
            filename=file_path.name,
            file_type="markdown",
            sections=sections
        )
    
    def _extract_entities(self, text: str) -> List[ExtractedEntity]:
        """Extract named entities from text using SpaCy."""
        if not self.nlp:
            return []
        
        doc = self.nlp(text)
        entities = []
        
        for ent in doc.ents:
            # Filter out common/low-value entities
            if ent.label_ in ['PERSON', 'ORG', 'GPE', 'PRODUCT', 'EVENT', 'MONEY', 'PERCENT']:
                entity = ExtractedEntity(
                    text=ent.text,
                    label=ent.label_,
                    confidence=getattr(ent, 'confidence', 0.9),  # SpaCy doesn't provide confidence by default
                    start_char=ent.start_char,
                    end_char=ent.end_char
                )
                entities.append(entity)
        
        return entities
    
    def _extract_table_text(self, table) -> str:
        """Extract text content from a docx table."""
        table_data = []
        for row in table.rows:
            row_data = []
            for cell in row.cells:
                row_data.append(cell.text.strip())
            table_data.append(" | ".join(row_data))
        
        return "\n".join(table_data)
    
    def _create_processed_document(
        self,
        filename: str,
        file_type: str,
        sections: List[DocumentSection],
        total_pages: Optional[int] = None,
        total_slides: Optional[int] = None
    ) -> ProcessedDocument:
        """Create ProcessedDocument with extracted content and analysis."""
        
        # Aggregate all entities
        all_entities = []
        all_content = []
        
        for section in sections:
            all_entities.extend(section.entities)
            all_content.append(section.content)
        
        full_text = "\n\n".join(all_content)
        
        # Extract key facts (simplified approach)
        key_facts = self._extract_key_facts(full_text)
        
        # Generate summary (simplified)
        summary = self._generate_summary(full_text)
        
        return ProcessedDocument(
            filename=filename,
            file_type=file_type,
            total_pages=total_pages,
            total_slides=total_slides,
            sections=sections,
            entities=all_entities,
            key_facts=key_facts,
            summary=summary,
            metadata={
                "total_sections": len(sections),
                "total_entities": len(all_entities),
                "word_count": len(full_text.split()),
                "character_count": len(full_text)
            }
        )
    
    def _extract_key_facts(self, text: str) -> List[str]:
        """Extract key facts from text using simple heuristics."""
        sentences = nltk.sent_tokenize(text)
        key_facts = []
        
        # Look for sentences with numbers, percentages, or important keywords
        fact_indicators = [
            r'\d+%',  # percentages
            r'\$\d+',  # money
            r'\d{4}',  # years
            r'increase|decrease|grow|decline',  # trends
            r'market|revenue|profit|cost',  # business terms
            r'strategy|goal|objective|target',  # strategic terms
        ]
        
        pattern = '|'.join(fact_indicators)
        
        for sentence in sentences:
            if re.search(pattern, sentence, re.IGNORECASE) and len(sentence.split()) > 5:
                key_facts.append(sentence.strip())
        
        return key_facts[:10]  # Limit to top 10 facts
    
    def _generate_summary(self, text: str, max_sentences: int = 3) -> str:
        """Generate a simple extractive summary."""
        sentences = nltk.sent_tokenize(text)
        
        if len(sentences) <= max_sentences:
            return text
        
        # Simple approach: take first, middle, and last sentences
        indices = [0, len(sentences) // 2, -1]
        summary_sentences = [sentences[i] for i in indices[:max_sentences]]
        
        return " ".join(summary_sentences)
    
    def process_multiple_files(self, file_paths: List[Union[str, Path]]) -> List[ProcessedDocument]:
        """Process multiple files and return list of processed documents."""
        processed_docs = []
        
        for file_path in file_paths:
            try:
                doc = self.process_file(file_path)
                processed_docs.append(doc)
                logger.info(f"Successfully processed: {file_path}")
            except Exception as e:
                logger.error(f"Failed to process {file_path}: {e}")
        
        return processed_docs